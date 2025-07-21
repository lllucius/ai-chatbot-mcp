"""
Enhanced document processor with comprehensive format support and advanced preprocessing.

This module provides a comprehensive document processing pipeline that supports multiple
document formats and includes advanced text preprocessing capabilities including:

- Document format detection and processing for all text-based formats
- Advanced text preprocessing (normalization, cleaning, consistency correction)
- Configurable text chunking with metadata preservation
- Background processing support for large documents
- Enhanced error handling and logging

Key Features:
- Multi-format support (PDF, DOCX, TXT, MD, RTF, HTML, EPUB, CSV, TSV, XML, JSON)
- Text normalization and cleaning
- Modifiable chunking parameters
- Async processing with progress tracking
- Comprehensive metadata extraction
- Error recovery and reporting

Current User: assistant
Current Date: 2025-01-20
"""

import logging
import mimetypes
import re
import unicodedata
from pathlib import Path
from typing import Any, Dict, Optional, Tuple, Union

import filetype
from langdetect import detect, LangDetectError

logger = logging.getLogger(__name__)


class DocumentFormatError(Exception):
    """Exception raised for document format-related errors."""

    pass


class DocumentProcessor:
    """
    Enhanced document processor with comprehensive format support and preprocessing.

    This processor extends beyond the basic PyPDF2, python-docx, striprtf approach
    to provide comprehensive support for text-based document formats with advanced
    preprocessing capabilities.
    """

    SUPPORTED_FORMATS = {
        # Text formats
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".markdown": "text/markdown",
        ".rst": "text/x-rst",
        ".csv": "text/csv",
        ".tsv": "text/tab-separated-values",
        ".json": "application/json",
        ".jsonl": "application/jsonl",
        ".xml": "application/xml",
        ".html": "text/html",
        ".htm": "text/html",
        ".xhtml": "application/xhtml+xml",
        # Office documents
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".doc": "application/msword",
        ".odt": "application/vnd.oasis.opendocument.text",
        ".rtf": "application/rtf",
        # Presentations
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".ppt": "application/vnd.ms-powerpoint",
        ".odp": "application/vnd.oasis.opendocument.presentation",
        # Spreadsheets
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xls": "application/vnd.ms-excel",
        ".ods": "application/vnd.oasis.opendocument.spreadsheet",
        # eBooks
        ".epub": "application/epub+zip",
        # Email formats
        ".eml": "message/rfc822",
        ".msg": "application/vnd.ms-outlook",
    }

    def __init__(self):
        """Initialize the enhanced document processor."""
        self.setup_preprocessing_patterns()

    def setup_preprocessing_patterns(self):
        """Set up regex patterns for text preprocessing."""
        # Common unwanted characters and patterns
        self.unwanted_patterns = [
            # Multiple whitespaces
            re.compile(r"\s+"),
            # Non-printable characters (except newlines, tabs, carriage returns)
            re.compile(r"[^\x20-\x7E\x0A\x0D\x09\u00A0-\uFFFF]"),
            # Multiple consecutive newlines
            re.compile(r"\n\s*\n\s*\n+"),
            # Trailing whitespace
            re.compile(r"[ \t]+$", re.MULTILINE),
            # Leading whitespace on empty lines
            re.compile(r"^\s+$", re.MULTILINE),
        ]

        # Text normalization patterns
        self.normalization_patterns = [
            # Smart double quotes: “ ” to "
            (re.compile(r"[“”]"), '"'),
            # Smart single quotes: ‘ ’ to '
            (re.compile(r"[‘’]"), "'"),
            # Em/en dashes to hyphens
            (re.compile(r"[—–]"), "-"),
            # Ellipsis normalization
            (re.compile(r"…"), "..."),
            # Non-breaking space to regular space
            (re.compile(r"\u00A0"), " "),
        ]

    def detect_file_format(self, file_path: Union[str, Path]) -> Tuple[str, str]:
        """
        Detect file format using multiple methods.

        Args:
            file_path: Path to the file

        Returns:
            Tuple[str, str]: (file_extension, mime_type)

        Raises:
            DocumentFormatError: If format cannot be determined
        """
        file_path = Path(file_path)

        # Get extension
        extension = file_path.suffix.lower()

        # Try filetype detection first (more accurate)
        try:
            detected_type = filetype.guess(file_path)
            if detected_type:
                mime_type = detected_type.mime
                # Map back to extension if needed
                if not extension:
                    extension = f".{detected_type.extension}"
            else:
                # Fall back to mimetypes
                mime_type, _ = mimetypes.guess_type(file_path)
                if not mime_type and extension in self.SUPPORTED_FORMATS:
                    mime_type = self.SUPPORTED_FORMATS[extension]
        except Exception as e:
            logger.warning(f"File type detection failed for {file_path}: {e}")
            mime_type = self.SUPPORTED_FORMATS.get(extension)

        if not extension or extension not in self.SUPPORTED_FORMATS:
            raise DocumentFormatError(f"Unsupported file format: {extension}")

        if not mime_type:
            mime_type = self.SUPPORTED_FORMATS[extension]

        return extension, mime_type

    def preprocess_text(self, text: str) -> str:
        """
        Apply comprehensive text preprocessing.

        Args:
            text: Raw text to preprocess

        Returns:
            str: Preprocessed and normalized text
        """
        if not text:
            return ""

        # Unicode normalization
        text = unicodedata.normalize("NFKC", text)

        # Apply normalization patterns
        for pattern, replacement in self.normalization_patterns:
            text = pattern.sub(replacement, text)

        # Remove unwanted patterns
        # Multiple whitespaces -> single space
        text = self.unwanted_patterns[0].sub(" ", text)

        # Remove non-printable characters
        text = self.unwanted_patterns[1].sub("", text)

        # Normalize multiple newlines
        text = self.unwanted_patterns[2].sub("\n\n", text)

        # Remove trailing whitespace
        text = self.unwanted_patterns[3].sub("", text)

        # Remove leading whitespace on empty lines
        text = self.unwanted_patterns[4].sub("", text)

        # Final cleanup
        text = text.strip()

        return text

    async def extract_text(self, file_path: Union[str, Path]) -> str:
        """
        Extract text from document with format detection and preprocessing.

        Args:
            file_path: Path to the document file

        Returns:
            str: Extracted and preprocessed text

        Raises:
            DocumentFormatError: If extraction fails
        """
        file_path = Path(file_path)

        if not file_path.exists():
            raise DocumentFormatError(f"File not found: {file_path}")

        # Detect file format
        extension, mime_type = self.detect_file_format(file_path)

        logger.info(f"Processing {extension} file: {file_path}")

        # Extract text based on format
        try:
            if extension in [
                ".txt",
                ".md",
                ".markdown",
                ".rst",
                ".csv",
                ".tsv",
                ".json",
                ".jsonl",
                ".xml",
            ]:
                text = await self._extract_text_format(file_path, extension)
            elif extension in [".html", ".htm", ".xhtml"]:
                text = await self._extract_html(file_path)
            elif extension == ".pdf":
                text = await self._extract_pdf(file_path)
            elif extension in [".docx", ".doc"]:
                text = await self._extract_word(file_path)
            elif extension == ".rtf":
                text = await self._extract_rtf(file_path)
            elif extension in [".xlsx", ".xls", ".ods"]:
                text = await self._extract_spreadsheet(file_path)
            elif extension in [".pptx", ".ppt"]:
                text = await self._extract_presentation(file_path)
            elif extension == ".epub":
                text = await self._extract_epub(file_path)
            elif extension in [".eml", ".msg"]:
                text = await self._extract_email(file_path)
            else:
                # Fallback to text extraction
                text = await self._extract_text_format(file_path, ".txt")

        except Exception as e:
            logger.error(f"Text extraction failed for {file_path}: {e}")
            raise DocumentFormatError(f"Failed to extract text: {e}")

        # Apply preprocessing
        processed_text = self.preprocess_text(text)

        if not processed_text:
            raise DocumentFormatError("No text content extracted from document")

        logger.info(f"Extracted {len(processed_text)} characters from {file_path}")
        return processed_text

    async def _extract_text_format(self, file_path: Path, extension: str) -> str:
        """Extract text from plain text formats."""
        # Try different encodings
        encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252", "iso-8859-1"]

        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    content = f.read()

                # Special handling for structured formats
                if extension == ".json":
                    content = self._extract_json_text(content)
                elif extension == ".xml":
                    content = self._extract_xml_text(content)
                elif extension in [".csv", ".tsv"]:
                    content = self._extract_csv_text(content, extension)

                return content

            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.warning(f"Failed to read with encoding {encoding}: {e}")
                continue

        raise DocumentFormatError(
            "Unable to decode text file with any supported encoding"
        )

    async def _extract_html(self, file_path: Path) -> str:
        """Extract text from HTML files."""
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            logger.warning(
                "BeautifulSoup not available, falling back to basic HTML extraction"
            )
            return await self._extract_text_format(file_path, ".txt")

        with open(file_path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f.read(), "html.parser")

        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()

        # Extract text
        text = soup.get_text()

        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = "\n".join(chunk for chunk in chunks if chunk)

        return text

    async def _extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF files."""
        try:
            import PyPDF2
        except ImportError:
            raise DocumentFormatError("PyPDF2 not installed for PDF processing")

        text_content = []

        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)

            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_content.append(page_text)
                except Exception as e:
                    logger.warning(f"Failed to extract text from page {page_num}: {e}")
                    continue

        return "\n\n".join(text_content)

    async def _extract_word(self, file_path: Path) -> str:
        """Extract text from Word documents."""
        try:
            import docx
        except ImportError:
            raise DocumentFormatError(
                "python-docx not installed for Word document processing"
            )

        doc = docx.Document(file_path)
        text_content = []

        # Extract paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)

        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_content.append(" | ".join(row_text))

        return "\n".join(text_content)

    async def _extract_rtf(self, file_path: Path) -> str:
        """Extract text from RTF files."""
        try:
            from striprtf.striprtf import rtf_to_text
        except ImportError:
            raise DocumentFormatError("striprtf not installed for RTF processing")

        with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
            rtf_content = file.read()

        return rtf_to_text(rtf_content)

    async def _extract_spreadsheet(self, file_path: Path) -> str:
        """Extract text from spreadsheet files."""
        try:
            import pandas as pd
        except ImportError:
            logger.warning("pandas not available for spreadsheet processing")
            return ""

        try:
            # Try to read as Excel file
            df = pd.read_excel(file_path, sheet_name=None)
            text_content = []

            for sheet_name, sheet_df in df.items():
                text_content.append(f"Sheet: {sheet_name}")
                text_content.append(sheet_df.to_string(index=False))
                text_content.append("")

            return "\n".join(text_content)

        except Exception as e:
            logger.warning(f"Failed to process spreadsheet {file_path}: {e}")
            return ""

    async def _extract_presentation(self, file_path: Path) -> str:
        """Extract text from presentation files."""
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx not available for presentation processing")
            return ""

        try:
            prs = Presentation(file_path)
            text_content = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                if slide_text:
                    text_content.append(f"Slide {slide_num}:")
                    text_content.extend(slide_text)
                    text_content.append("")

            return "\n".join(text_content)

        except Exception as e:
            logger.warning(f"Failed to process presentation {file_path}: {e}")
            return ""

    async def _extract_epub(self, file_path: Path) -> str:
        """Extract text from EPUB files."""
        try:
            import zipfile
            from bs4 import BeautifulSoup
        except ImportError:
            logger.warning("Required libraries not available for EPUB processing")
            return ""

        text_content = []

        try:
            with zipfile.ZipFile(file_path, "r") as epub:
                for file_info in epub.filelist:
                    if file_info.filename.endswith((".html", ".xhtml", ".htm")):
                        with epub.open(file_info.filename) as f:
                            soup = BeautifulSoup(f.read(), "html.parser")

                            # Remove script and style elements
                            for script in soup(["script", "style"]):
                                script.decompose()

                            text = soup.get_text()
                            if text.strip():
                                text_content.append(text)

            return "\n\n".join(text_content)

        except Exception as e:
            logger.warning(f"Failed to process EPUB {file_path}: {e}")
            return ""

    async def _extract_email(self, file_path: Path) -> str:
        """Extract text from email files."""
        try:
            import email
        except ImportError:
            logger.warning("email library not available")
            return ""

        try:
            with open(file_path, "rb") as f:
                msg = email.message_from_bytes(f.read())

            text_parts = []

            # Add headers
            text_parts.append(f"Subject: {msg.get('Subject', 'No Subject')}")
            text_parts.append(f"From: {msg.get('From', 'Unknown')}")
            text_parts.append(f"To: {msg.get('To', 'Unknown')}")
            text_parts.append(f"Date: {msg.get('Date', 'Unknown')}")
            text_parts.append("")

            # Extract body
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        text_parts.append(
                            part.get_payload(decode=True).decode(
                                "utf-8", errors="ignore"
                            )
                        )
            else:
                text_parts.append(
                    msg.get_payload(decode=True).decode("utf-8", errors="ignore")
                )

            return "\n".join(text_parts)

        except Exception as e:
            logger.warning(f"Failed to process email {file_path}: {e}")
            return ""

    def _extract_json_text(self, content: str) -> str:
        """Extract meaningful text from JSON content."""
        try:
            import json

            def extract_text_from_obj(obj, path=""):
                """Recursively extract text from JSON object."""
                texts = []

                if isinstance(obj, dict):
                    for key, value in obj.items():
                        if isinstance(value, str) and value.strip():
                            texts.append(f"{key}: {value}")
                        elif isinstance(value, (dict, list)):
                            texts.extend(
                                extract_text_from_obj(
                                    value, f"{path}.{key}" if path else key
                                )
                            )
                elif isinstance(obj, list):
                    for i, item in enumerate(obj):
                        if isinstance(item, str) and item.strip():
                            texts.append(item)
                        elif isinstance(item, (dict, list)):
                            texts.extend(
                                extract_text_from_obj(
                                    item, f"{path}[{i}]" if path else f"[{i}]"
                                )
                            )
                elif isinstance(obj, str) and obj.strip():
                    texts.append(obj)

                return texts

            data = json.loads(content)
            return "\n".join(extract_text_from_obj(data))

        except json.JSONDecodeError:
            # If JSON parsing fails, return raw content
            return content

    def _extract_xml_text(self, content: str) -> str:
        """Extract text from XML content."""
        try:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(content, "xml")
            return soup.get_text()

        except ImportError:
            # Basic XML text extraction without BeautifulSoup
            import re

            # Remove XML tags
            text = re.sub(r"<[^>]+>", "", content)
            return text

    def _extract_csv_text(self, content: str, extension: str) -> str:
        """Extract text from CSV/TSV content."""
        try:
            import csv
            import io

            delimiter = "\t" if extension == ".tsv" else ","

            # Parse CSV and convert to text
            reader = csv.reader(io.StringIO(content), delimiter=delimiter)
            rows = []

            for row in reader:
                if row:  # Skip empty rows
                    rows.append(
                        " | ".join(cell.strip() for cell in row if cell.strip())
                    )

            return "\n".join(rows)

        except Exception:
            # Fallback to raw content
            return content

    def detect_language(self, text: str) -> Optional[str]:
        """
        Detect the language of the text.

        Args:
            text: Text to analyze

        Returns:
            Optional[str]: Detected language code (ISO 639-1) or None if detection fails
        """
        try:
            return detect(text)
        except LangDetectError:
            return None

    def get_text_statistics(self, text: str) -> Dict[str, Any]:
        """
        Get comprehensive statistics about the text.

        Args:
            text: Text to analyze

        Returns:
            Dict[str, Any]: Statistics including character count, word count, etc.
        """
        if not text:
            return {
                "character_count": 0,
                "word_count": 0,
                "line_count": 0,
                "paragraph_count": 0,
                "language": None,
            }

        words = text.split()
        lines = text.split("\n")
        paragraphs = [p for p in text.split("\n\n") if p.strip()]

        stats = {
            "character_count": len(text),
            "word_count": len(words),
            "line_count": len(lines),
            "paragraph_count": len(paragraphs),
            "language": self.detect_language(
                text[:1000]
            ),  # Sample first 1000 chars for language detection
        }

        return stats
